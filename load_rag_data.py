"""Script to load Google Drive data into RAG vector database."""
import sys
import requests
import io
from typing import Dict
from config import DOCTOR_PROFILES, PATIENT_DATA_LINK
import PyPDF2
from docx import Document
from pptx import Presentation


def extract_file_id_from_url(url: str) -> str:
    """Extract file ID from Google Drive share URL."""
    if not url:
        return None

    if "/d/" in url:
        return url.split("/d/")[1].split("/")[0]
    elif "?id=" in url:
        return url.split("?id=")[1].split("&")[0]

    return None


def download_from_share_link(url: str) -> bytes:
    """Download file from Google Drive shared link."""
    try:
        file_id = extract_file_id_from_url(url)
        if not file_id:
            print(f"[ERROR] Could not extract file ID from URL: {url}")
            return None

        # Use export URL for direct download
        download_url = f"https://drive.google.com/uc?export=download&id={file_id}"

        print(f"   Downloading: {file_id[:10]}...")
        response = requests.get(download_url, timeout=30)

        if response.status_code == 200:
            return response.content
        else:
            print(f"   [ERROR] Download failed (status {response.status_code})")
            return None

    except Exception as e:
        print(f"   [ERROR] Error downloading file: {e}")
        return None


def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """Extract text from PDF bytes."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text if text.strip() else None
    except Exception as e:
        return None


def extract_text_from_docx(docx_bytes: bytes) -> str:
    """Extract text from DOCX bytes."""
    try:
        doc = Document(io.BytesIO(docx_bytes))
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text if text.strip() else None
    except Exception as e:
        return None


def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """Extract text from PPTX bytes."""
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if text.strip() else None
    except Exception as e:
        return None


def extract_text_from_txt(txt_bytes: bytes) -> str:
    """Extract text from TXT bytes."""
    try:
        return txt_bytes.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        return None


def extract_text_from_csv(csv_bytes: bytes) -> str:
    """Extract text from CSV bytes."""
    try:
        return csv_bytes.decode("utf-8", errors="ignore").strip()
    except Exception as e:
        return None


def extract_file_content(file_bytes: bytes, file_name: str) -> str:
    """Extract content from various file types.

    Tries multiple formats if extension is unknown, since Google Drive
    URLs don't always preserve file extensions.
    """
    if not file_bytes:
        return None

    file_name_lower = file_name.lower()

    # Try based on explicit extension first
    if file_name_lower.endswith(".txt"):
        result = extract_text_from_txt(file_bytes)
        if result:
            return result
    elif file_name_lower.endswith(".csv"):
        result = extract_text_from_csv(file_bytes)
        if result:
            return result
    elif file_name_lower.endswith(".pdf"):
        result = extract_text_from_pdf(file_bytes)
        if result:
            return result
    elif file_name_lower.endswith(".docx"):
        result = extract_text_from_docx(file_bytes)
        if result:
            return result
    elif file_name_lower.endswith(".pptx"):
        result = extract_text_from_pptx(file_bytes)
        if result:
            return result

    # If no extension or unknown, try all formats
    # Try text first (most common)
    result = extract_text_from_txt(file_bytes)
    if result:
        return result

    # Try CSV
    result = extract_text_from_csv(file_bytes)
    if result:
        return result

    # Try other formats
    result = extract_text_from_docx(file_bytes)
    if result:
        return result

    result = extract_text_from_pptx(file_bytes)
    if result:
        return result

    result = extract_text_from_pdf(file_bytes)
    if result:
        return result

    return None


def load_doctor_profiles() -> Dict[str, str]:
    """Load doctor profiles from Google Drive shared links."""
    profiles = {}

    for doctor_name, drive_link in DOCTOR_PROFILES.items():
        if not drive_link:
            print(f"   [SKIP]  Skipping {doctor_name} (no link configured)")
            continue

        print(f"   Loading {doctor_name}...")
        file_bytes = download_from_share_link(drive_link)

        if file_bytes:
            # Try to extract without assuming file type
            content = extract_file_content(file_bytes, doctor_name)
            if content:
                profiles[doctor_name] = content
                print(f"   [OK] {doctor_name}")
            else:
                print(f"   [ERROR] Could not extract text from {doctor_name}")
        else:
            print(f"   [ERROR] Could not download {doctor_name}")

    return profiles


def load_patient_data() -> str:
    """Load patient data from Google Drive shared link."""
    if not PATIENT_DATA_LINK:
        print("   [SKIP]  Patient data link not configured")
        return None

    print("   Loading patient data...")
    file_bytes = download_from_share_link(PATIENT_DATA_LINK)

    if file_bytes:
        # Try to extract without assuming file type (could be CSV, TXT, etc.)
        content = extract_file_content(file_bytes, "patient_data")
        if content:
            print("   [OK] Patient data loaded")
            return content
        else:
            print("   [ERROR] Could not extract text from patient data")
    else:
        print("   [ERROR] Could not download patient data")

    return None


def load_rag_data(force_reload: bool = False):
    """Load Google Drive data into RAG vector database."""
    print("\n[DATA] Loading RAG data from Google Drive...\n")

    # Load doctor profiles
    print("[1]  Loading doctor profiles...")
    doctor_profiles = load_doctor_profiles()

    if not doctor_profiles:
        print("[WARN]  No doctor profiles loaded")
        return

    # Load patient data
    print("\n[2]  Loading patient data...")
    patient_data = load_patient_data()

    # Prepare data for vector DB
    print("\n[3]  Preparing data for vector database...")
    google_data = {
        "doctor_profiles": "\n\n".join(
            [f"## {name}\n{content}" for name, content in doctor_profiles.items()]
        )
    }

    if patient_data:
        google_data["patient_data"] = patient_data

    # Initialize vector database (lazy import)
    print("[4]  Initializing vector database...")
    try:
        from rag_vector_db import initialize_rag_db
        rag_db = initialize_rag_db()
    except ImportError as e:
        print(f"[ERROR] Error: RAG vector DB not available: {e}")
        print("   Run: pip install -r requirements.txt")
        return
    except Exception as e:
        print(f"[ERROR] Error initializing vector DB: {e}")
        return

    # Clear existing data if force reload
    if force_reload:
        print(" Force reload enabled - clearing existing data...")
        try:
            rag_db.clear_db()
        except Exception as e:
            print(f"[WARN]  Could not clear DB: {e}")

    # Add documents to vector store
    print("5  Loading documents into vector database...")
    try:
        rag_db.add_documents(google_data)
    except Exception as e:
        print(f"[ERROR] Error adding documents: {e}")
        return

    # Display statistics
    print("\n6  Database Statistics:")
    try:
        stats = rag_db.get_db_stats()
        for key, value in stats.items():
            print(f"   {key}: {value}")
    except Exception as e:
        print(f"[WARN]  Could not get stats: {e}")

    print("\n[OK] RAG data loading complete!\n")


def main():
    """Main entry point."""
    force_reload = "--force" in sys.argv or "-f" in sys.argv
    load_rag_data(force_reload=force_reload)


if __name__ == "__main__":
    main()
